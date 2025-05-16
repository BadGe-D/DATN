from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Helper function to add a slide with title and content
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Create a presentation
prs = Presentation()

# Slide 1: Giới thiệu đề tài
add_slide(prs, "Giới thiệu đề tài", 
          "Nghiên cứu và xây dựng hệ thống hỗ trợ di chuyển cho người khiếm thị và người có thị lực kém.\n"
          "Mục tiêu là sinh mô tả đường đi và chuyển thành âm thanh để hướng dẫn người dùng.")

# Slide 2: Tổng quan hệ thống
add_slide(prs, "Tổng quan hệ thống", 
          "- Bộ mã hóa ảnh (InceptionV3)\n- Bộ giải mã văn bản (LSTM)\n- Trình tạo câu (Greedy/Beam Search)\n- TTS chuyển văn bản thành âm thanh.")

# Slide 3: Kiến trúc Merging Architecture
add_slide(prs, "Kiến trúc Merging Architecture", 
          "Thông tin ảnh được giữ riêng và trộn sau khi LSTM mã hóa tiền tố văn bản.\n"
          "Kết hợp ảnh và ngữ cảnh tại lớp ngoài RNN.")

# Slide 4: Bộ giải mã và sinh câu
add_slide(prs, "Bộ giải mã và sinh câu", 
          "Sử dụng LSTM. Đầu vào là đặc trưng ảnh + từ trước đó.\nSinh từng từ theo trình tự đến khi gặp từ 'end'.")

# Slide 5: Thuật toán Greedy Search
add_slide(prs, "Thuật toán Greedy Search", 
          "Chọn từ có xác suất cao nhất ở mỗi bước.\nƯu: đơn giản, nhanh.\nNhược: có thể bỏ qua các phương án tốt hơn.")

# Slide 6: Thuật toán Beam Search
add_slide(prs, "Thuật toán Beam Search", 
          "Duy trì k lựa chọn tốt nhất ở mỗi bước.\nƯu: câu mượt hơn.\nNhược: tốn tài nguyên tính toán hơn.")

# Slide 7: Quy trình hệ thống
add_slide(prs, "Quy trình hoạt động toàn hệ thống", 
          "1. Nhận ảnh\n2. Trích xuất đặc trưng\n3. Sinh mô tả\n4. TTS\n5. Phát âm thanh cho người dùng")

# Slide 8: Kết quả và đánh giá
add_slide(prs, "Kết quả và đánh giá", 
          "Greedy: nhanh, dễ bị lặp.\nBeam: mô tả ngữ nghĩa tốt hơn.\nBLEU Score được dùng để đánh giá.")

# Slide 9: Kết luận và hướng phát triển
add_slide(prs, "Kết luận và hướng phát triển", 
          "- Tích hợp GPS\n- Nhận dạng vật thể nâng cao\n- Tối ưu cho thiết bị di động\n- Hỗ trợ điều kiện môi trường đa dạng")

# Save the presentation
pptx_path = "TrinhBay_Merging_Greedy_Beam_Search.pptx"
prs.save(pptx_path)

pptx_path
